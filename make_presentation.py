"""Gera a apresentacao do case Cloud Humans em PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta de cores ──────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # azul escuro (fundo titulo)
MID_BG    = RGBColor(0x16, 0x21, 0x3E)   # azul medio
ACCENT    = RGBColor(0x0F, 0x3A, 0x60)   # azul destaque
TEAL      = RGBColor(0x00, 0xB4, 0xD8)   # ciano/teal
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TXT = RGBColor(0xD0, 0xE8, 0xF2)
YELLOW    = RGBColor(0xFF, 0xD1, 0x66)
GREEN     = RGBColor(0x06, 0xD6, 0xA0)
ORANGE    = RGBColor(0xFF, 0x74, 0x00)
GRAY_BG   = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TXT  = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout em branco


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def txt_lines(slide, lines, l, t, w, h, size=16, color=WHITE,
              bold_first=False, bullet=False, leading_space=None):
    """Adiciona multiplas linhas em um unico textbox."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if leading_space and not first:
            p.space_before = Pt(leading_space)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if bullet:
            run.text = ("• " if not line.startswith("•") else "") + line
        else:
            run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold_first and (line == lines[0])
        run.font.color.rgb = color
    return box


def dark_slide(title_txt, subtitle_txt=None):
    """Slide de secao com fundo escuro."""
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    add_rect(slide, 0, 0, 0.15, 7.5, TEAL)          # barra lateral esquerda
    add_rect(slide, 0.15, 3.3, 13.18, 0.06, TEAL)   # linha divisoria

    txt(slide, title_txt, 0.5, 2.5, 12.5, 1.2, size=36, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle_txt:
        txt(slide, subtitle_txt, 0.5, 3.6, 12.5, 0.7, size=20,
            color=LIGHT_TXT, align=PP_ALIGN.CENTER, italic=True)
    return slide


def content_slide(title_txt, bg_left=True):
    """Slide de conteudo com fundo claro."""
    slide = prs.slides.add_slide(BLANK)
    # fundo
    add_rect(slide, 0, 0, 13.33, 7.5, GRAY_BG)
    # header bar
    add_rect(slide, 0, 0, 13.33, 1.1, DARK_BG)
    add_rect(slide, 0, 0, 0.15, 7.5, TEAL)
    txt(slide, title_txt, 0.3, 0.15, 12.8, 0.85, size=22, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT)
    return slide


def kpi_box(slide, label, value, l, t, w=2.8, h=1.6, bg=ACCENT):
    add_rect(slide, l, t, w, h, bg)
    txt(slide, value, l+0.1, t+0.15, w-0.2, 0.75, size=28, bold=True,
        color=YELLOW, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.1, t+0.85, w-0.2, 0.6, size=13,
        color=LIGHT_TXT, align=PP_ALIGN.CENTER)


def step_box(slide, number, title, desc, l, t, w=3.8, h=2.0, bg=ACCENT):
    add_rect(slide, l, t, w, h, bg)
    add_rect(slide, l, t, w, 0.45, TEAL)
    txt(slide, number, l+0.1, t+0.05, 0.6, 0.38, size=18, bold=True,
        color=DARK_BG)
    txt(slide, title, l+0.65, t+0.05, w-0.75, 0.38, size=14, bold=True,
        color=DARK_BG)
    txt(slide, desc, l+0.15, t+0.55, w-0.3, 1.3, size=12,
        color=LIGHT_TXT, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(slide, 0, 0, 0.15, 7.5, TEAL)
add_rect(slide, 0.15, 6.8, 13.18, 0.08, TEAL)
# decoracao geometrica
add_rect(slide, 9.5, 0, 3.83, 7.5, ACCENT)
add_rect(slide, 9.5, 0, 0.06, 7.5, TEAL)

txt(slide, "Cloud Humans", 0.5, 1.2, 8.7, 0.9, size=16, color=TEAL, italic=True)
txt(slide, "Data Engineering\nCase", 0.5, 1.95, 8.7, 2.0, size=50, bold=True,
    color=WHITE, align=PP_ALIGN.LEFT)
txt(slide, "Pipeline local em DuckDB para calcular métricas\nde tempo de serviço a partir dos CSVs crus",
    0.5, 4.1, 8.5, 1.2, size=18, color=LIGHT_TXT)
txt(slide, "Leonardo Correia  •  Junho 2026", 0.5, 6.5, 8.5, 0.6,
    size=14, color=TEAL, italic=True)

txt(slide, "DuckDB", 9.8, 1.5, 3.2, 0.6, size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(slide, "SQL Analítico", 9.8, 2.1, 3.2, 0.5, size=14, color=LIGHT_TXT, align=PP_ALIGN.CENTER)
add_rect(slide, 10.1, 2.65, 2.6, 0.06, TEAL)
txt(slide, "Python", 9.8, 2.8, 3.2, 0.6, size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(slide, "Orquestração local", 9.8, 3.4, 3.2, 0.5, size=14, color=LIGHT_TXT, align=PP_ALIGN.CENTER)
add_rect(slide, 10.1, 3.95, 2.6, 0.06, TEAL)
txt(slide, "6 Passos", 9.8, 4.1, 3.2, 0.6, size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(slide, "Da ingestão ao relatório", 9.8, 4.7, 3.2, 0.5, size=14, color=LIGHT_TXT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTO / PROBLEMA
# ════════════════════════════════════════════════════════════════════════════
slide = content_slide("O Problema: Por Que Este Case Existe")
add_rect(slide, 0.3, 1.25, 12.7, 5.9, ACCENT)
add_rect(slide, 0.3, 1.25, 0.08, 5.9, ORANGE)

txt(slide, "A Cloud Humans enfrenta 3 dores no dashboard de atendimento:",
    0.55, 1.35, 12.3, 0.55, size=16, bold=True, color=YELLOW)

dores = [
    ("1  Lentidão / Performance",
     "O dashboard dispara queries pesadas em tempo real: joins, janelas temporais, deduplicações\n"
     "e agregações sobre dados operacionais brutos a cada interação do usuário."),
    ("2  Falta de Visibilidade",
     "Entidades analíticas críticas não existem pré-calculadas: jornada da conversa, tempo em\n"
     "fila, tempo humano vs IA, tempo após escalonamento."),
    ("3  Regras de Negócio Mal Definidas",
     "Cada gráfico implementa sua própria query → a mesma métrica pode ter definições\n"
     "diferentes em dashboards distintos, gerando divergência de números."),
]

y = 2.0
for titulo, desc in dores:
    add_rect(slide, 0.45, y, 12.4, 1.4, MID_BG)
    txt(slide, titulo, 0.6, y+0.1, 12.0, 0.45, size=15, bold=True, color=TEAL)
    txt(slide, desc, 0.6, y+0.52, 12.0, 0.78, size=13, color=LIGHT_TXT)
    y += 1.55

txt(slide, "Raiz do problema: o BI consulta dados brutos e recalcula tudo no momento da visualização.",
    0.45, 6.8, 12.4, 0.4, size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ABORDAGEM / VISAO GERAL
# ════════════════════════════════════════════════════════════════════════════
slide = content_slide("Abordagem: Pipeline em Camadas com DuckDB")

txt(slide, "Simula localmente a arquitetura em camadas que resolveria o problema em produção:",
    0.4, 1.2, 12.5, 0.45, size=15, color=DARK_TXT)

steps = [
    ("CSVs\nCrus", "4 arquivos\nde entrada", ORANGE),
    ("Raw\nTables", "Ingestão\ndireta", RGBColor(0x02, 0x96, 0xB3)),
    ("Staging\nLimpo", "Tipagem\n+ dedup", TEAL),
    ("Jornada\nGranular", "Intervalos\npor estado", GREEN),
    ("Métricas\nFinais", "Por conversa", RGBColor(0x9C, 0x27, 0xB0)),
    ("Relatório\nVisual", "HTML + CSV", YELLOW),
]

x = 0.35
for label, sub, color in steps:
    add_rect(slide, x, 1.9, 2.0, 1.3, color)
    txt(slide, label, x, 1.95, 2.0, 0.8, size=15, bold=True,
        color=DARK_BG if color in (YELLOW, TEAL, GREEN) else WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sub, x, 2.7, 2.0, 0.42, size=11,
        color=DARK_BG if color in (YELLOW, TEAL, GREEN) else WHITE, align=PP_ALIGN.CENTER)
    if x < 0.35 + 5 * 2.15:
        txt(slide, "→", x+2.0, 2.3, 0.25, 0.5, size=20, bold=True, color=DARK_TXT)
    x += 2.15

# Arquivos envolvidos
txt(slide, "Arquivos principais:", 0.4, 3.45, 12.5, 0.4, size=14, bold=True, color=DARK_TXT)

files = [
    ("duck_init.py", "Cria tabelas raw no DuckDB"),
    ("sql/02_staging.sql", "Limpeza, tipagem e deduplicação"),
    ("sql/03_journey.sql", "Reconstrói jornada por intervalos"),
    ("sql/04_service_metrics.sql", "Métricas por conversa"),
    ("sql/05_human_escalation_metrics.sql", "Escalonamento humano"),
    ("sql/06_stakeholder_report.sql + run_stakeholder_report.py", "Relatório HTML + CSVs"),
]

x_col = [0.4, 6.9]
y_base = 3.9
for i, (fname, desc) in enumerate(files):
    col = i % 2
    row = i // 2
    xl = x_col[col]
    yl = y_base + row * 0.72
    add_rect(slide, xl, yl, 6.1, 0.58, MID_BG)
    txt(slide, fname, xl+0.1, yl+0.05, 3.0, 0.48, size=12, bold=True, color=TEAL)
    txt(slide, desc, xl+3.1, yl+0.1, 2.9, 0.38, size=12, color=LIGHT_TXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PASSO 1: INGESTAO
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Passo 1 — Ingestão Bruta", "duck_init.py  ·  tabelas raw no DuckDB")

boxes = [
    ("accounts.csv", "Contas (clientes)\ncadastradas no sistema"),
    ("agents.csv", "Agentes (humanos e IA)\nque atendem conversas"),
    ("conversations.csv", "Cada conversa:\ncriação, resolução, inbox"),
    ("conversation_events.csv", "Eventos granulares:\natribuições, resoluções, bots"),
]

x = 0.6
for title, desc in boxes:
    add_rect(slide, x, 2.0, 2.7, 2.2, ACCENT)
    add_rect(slide, x, 2.0, 2.7, 0.42, TEAL)
    txt(slide, title, x+0.1, 2.05, 2.5, 0.35, size=13, bold=True, color=DARK_BG)
    txt(slide, desc, x+0.1, 2.5, 2.5, 1.55, size=13, color=LIGHT_TXT)
    x += 2.9

txt(slide, "O script lê cada CSV com pandas e cria a tabela correspondente via DuckDB (CREATE OR REPLACE TABLE).\n"
    "Nenhum dado é filtrado nesta etapa — todos os registros brutos são preservados para auditoria.",
    0.5, 4.55, 12.3, 1.1, size=14, color=LIGHT_TXT)

add_rect(slide, 0.5, 5.8, 12.3, 0.55, MID_BG)
txt(slide, "Resultado: 4 tabelas raw no arquivo case.duckdb prontas para a camada de staging",
    0.65, 5.87, 12.0, 0.4, size=14, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PASSO 2: STAGING
# ════════════════════════════════════════════════════════════════════════════
slide = content_slide("Passo 2 — Staging: Limpeza, Tipagem e Qualidade")

col_l = [0.35, 6.85]
col_w = 6.0

# Coluna esquerda — o que faz
add_rect(slide, 0.35, 1.25, 6.0, 5.9, ACCENT)
txt(slide, "O que o staging faz", 0.5, 1.3, 5.7, 0.45, size=15, bold=True, color=TEAL)

items_l = [
    "SELECT DISTINCT para remover duplicatas exatas",
    "Cast explícito de tipos (timestamps UTC, inteiros)",
    "TRIM em strings para remover espaços",
    "Chaves compostas únicas: instance + account_id\n  + conversation_id / event_id",
    "Fallback de created_at quando evento\n  conversation_created está ausente",
    "View data_quality_checks expõe anomalias",
]
y = 1.85
for item in items_l:
    txt(slide, "✓  " + item, 0.5, y, 5.7, 0.62, size=13, color=LIGHT_TXT)
    y += 0.63

# Coluna direita — qualidade dos dados
add_rect(slide, 6.85, 1.25, 6.1, 5.9, MID_BG)
txt(slide, "Imperfeições encontradas e tratadas", 7.0, 1.3, 5.8, 0.45, size=15, bold=True, color=TEAL)

issues = [
    ("40", "duplicatas exatas em conversations"),
    ("120", "duplicatas exatas em conversation_events"),
    ("394", "conversas sem evento conversation_created\n→ fallback para conversations.created_at"),
    ("15", "conversas sem eventos granulares"),
    ("257", "conversas sem evento de resolução"),
    ("335", "escaladas p/ humano sem first_human_reply"),
]
y = 1.85
for num, desc in issues:
    add_rect(slide, 7.0, y, 1.0, 0.58, ORANGE)
    txt(slide, num, 7.0, y+0.08, 1.0, 0.42, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, desc, 8.1, y+0.06, 4.6, 0.52, size=12, color=LIGHT_TXT)
    y += 0.72


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PASSO 3: JORNADA
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Passo 3 — Jornada Granular das Conversas", "sql/03_journey.sql")

txt(slide, "Reconstrói a linha do tempo de cada conversa em intervalos classificados por estado:",
    0.5, 1.75, 12.3, 0.5, size=15, color=LIGHT_TXT)

estados = [
    ("queue", "Sem agente atribuído\nConversa aguardando", ORANGE),
    ("ai_service", "Bot/IA atendendo\no cliente", TEAL),
    ("human_service", "Agente humano\natendendo", GREEN),
    ("resolved", "Conversa resolvida\n(pode reabrir)", RGBColor(0x9C, 0x27, 0xB0)),
    ("terminal", "Encerrada sem\nresolução formal", RGBColor(0x6B, 0x73, 0x80)),
]

x = 0.5
for estado, desc, color in estados:
    add_rect(slide, x, 2.45, 2.35, 1.7, color)
    txt(slide, estado, x+0.1, 2.5, 2.15, 0.5, size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    txt(slide, desc, x+0.1, 3.0, 2.15, 0.95, size=12, color=DARK_BG, align=PP_ALIGN.CENTER)
    x += 2.52

txt(slide, "Tabelas geradas:", 0.5, 4.35, 12.3, 0.4, size=14, bold=True, color=TEAL)

tabelas = [
    "analytics.conversation_event_timeline",
    "analytics.conversation_journey_intervals",
    "analytics.conversation_journey_event_summary",
    "analytics.conversation_journey_summary",
    "analytics.journey_quality_checks",
]
x = 0.5
for tab in tabelas:
    add_rect(slide, x, 4.85, 2.45, 0.52, MID_BG)
    txt(slide, tab.replace("analytics.", ""), x+0.1, 4.9, 2.25, 0.42, size=10, color=LIGHT_TXT)
    x += 2.55

txt(slide, "Cada intervalo tem: estado, duração em minutos, agente responsável (humano ou IA)",
    0.5, 5.5, 12.3, 0.45, size=13, color=LIGHT_TXT, italic=True)

add_rect(slide, 0.5, 6.1, 12.3, 0.55, ACCENT)
txt(slide, "Decisão: usamos staging.stg_conversations.created_at_utc como fallback quando "
    "conversation_created está ausente",
    0.65, 6.17, 12.0, 0.4, size=13, color=YELLOW, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PASSO 4: METRICAS DE SERVICO
# ════════════════════════════════════════════════════════════════════════════
slide = content_slide("Passo 4 — Métricas de Serviço por Conversa")

txt(slide, "sql/04_service_metrics.sql  →  analytics.service_time_metrics",
    0.4, 1.2, 12.5, 0.42, size=14, bold=True, color=ACCENT)

metricas = [
    ("resolution_time_minutes", "Criação → última resolução (ciclo completo para conversas reabertas)"),
    ("first_resolution_time_minutes", "Criação → primeira resolução (análise operacional)"),
    ("queue_minutes", "Total de tempo sem agente atribuído"),
    ("service_minutes", "Total de tempo com agente atribuído"),
    ("ai_service_minutes", "Tempo em atendimento por IA / bot"),
    ("human_service_minutes", "Tempo em atendimento por agente humano"),
    ("first_reply_time_minutes", "Criação → primeira resposta (qualquer agente)"),
    ("first_human_reply_time_minutes", "Criação → primeira resposta humana"),
    ("human_wait_until_first_reply_minutes", "Primeira atribuição humana → primeira resposta humana"),
]

y = 1.75
for i, (metric, desc) in enumerate(metricas):
    bg = ACCENT if i % 2 == 0 else MID_BG
    add_rect(slide, 0.35, y, 12.6, 0.5, bg)
    txt(slide, metric, 0.5, y+0.06, 5.0, 0.38, size=12, bold=True, color=TEAL)
    txt(slide, desc, 5.55, y+0.07, 7.2, 0.38, size=12, color=LIGHT_TXT)
    y += 0.52

txt(slide, "Também gera agregações: by_account, by_inbox, by_team, by_resolution_level",
    0.35, 6.55, 12.6, 0.45, size=13, bold=True, color=DARK_TXT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PASSO 5: ESCALONAMENTO HUMANO
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Passo 5 — Métricas de Escalonamento Humano", "sql/05_human_escalation_metrics.sql")

txt(slide, "Foco: o momento em que a IA transfere o atendimento para um humano.",
    0.5, 1.75, 12.3, 0.45, size=15, color=LIGHT_TXT)

add_rect(slide, 0.5, 2.3, 5.8, 3.3, ACCENT)
txt(slide, "Lógica Central", 0.65, 2.35, 5.5, 0.45, size=15, bold=True, color=TEAL)
items = [
    "Identifica a primeira atribuição a agente humano",
    "Registra a primeira resposta humana (via conversations)",
    "Calcula: first_human_assigned_at → first_human_reply",
    "Agrega por agente (instance + user_id) e por time",
    "Exposição em journey_quality_checks para casos sem reply",
]
y = 2.9
for item in items:
    txt(slide, "▸  " + item, 0.65, y, 5.5, 0.45, size=13, color=LIGHT_TXT)
    y += 0.48

add_rect(slide, 6.6, 2.3, 6.3, 3.3, MID_BG)
txt(slide, "Tabelas Geradas", 6.75, 2.35, 6.0, 0.45, size=15, bold=True, color=TEAL)
tabs = [
    ("analytics.first_human_assignment", "Primeira atribuição por conversa"),
    ("analytics.human_escalation_metrics", "Métrica por conversa"),
    ("analytics.human_escalation_by_agent", "Agregação por agente"),
    ("analytics.human_escalation_by_team", "Agregação por time"),
    ("analytics.human_escalation_quality_checks", "Checks de qualidade"),
]
y = 2.9
for tab, desc in tabs:
    txt(slide, tab.replace("analytics.", ""), 6.75, y, 6.0, 0.28, size=12, bold=True, color=YELLOW)
    txt(slide, desc, 6.75, y+0.28, 6.0, 0.22, size=11, color=LIGHT_TXT)
    y += 0.6

add_rect(slide, 0.5, 5.75, 12.3, 0.55, ORANGE)
txt(slide, "Nota: 335 conversas escaladas para humano não possuem first_human_reply_created_at "
    "→ flag de qualidade",
    0.65, 5.82, 12.0, 0.4, size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PASSO 6: RELATORIO
# ════════════════════════════════════════════════════════════════════════════
slide = content_slide("Passo 6 — Relatório Visual para Stakeholders")

add_rect(slide, 0.35, 1.25, 6.0, 5.9, ACCENT)
txt(slide, "O que é gerado", 0.5, 1.3, 5.7, 0.45, size=15, bold=True, color=TEAL)
outputs = [
    ("reports/service_metrics_report.html", "Relatório HTML completo com\ntabelas e KPIs formatados"),
    ("outputs/report_kpis.csv", "KPIs globais consolidados"),
    ("outputs/report_resolution_level.csv", "Métricas por nível (N1/N2)"),
    ("outputs/report_account_ranking.csv", "Ranking de contas por resolução"),
    ("outputs/report_human_wait_by_team.csv", "Espera humana por time"),
    ("outputs/report_human_wait_by_agent.csv", "Espera humana por agente"),
    ("outputs/report_key_insights.csv", "Insights automáticos extraídos"),
]
y = 1.85
for fname, desc in outputs:
    add_rect(slide, 0.45, y, 5.7, 0.72, MID_BG)
    txt(slide, fname, 0.55, y+0.04, 5.5, 0.32, size=11, bold=True, color=YELLOW)
    txt(slide, desc, 0.55, y+0.36, 5.5, 0.28, size=11, color=LIGHT_TXT)
    y += 0.77

add_rect(slide, 6.55, 1.25, 6.4, 5.9, MID_BG)
txt(slide, "Como executar", 6.7, 1.3, 6.1, 0.45, size=15, bold=True, color=TEAL)

txt(slide, "Pipeline completo (um comando):", 6.7, 1.85, 6.1, 0.38, size=13, color=LIGHT_TXT)
add_rect(slide, 6.7, 2.28, 6.1, 0.52, DARK_BG)
txt(slide, "python run_all.py", 6.85, 2.33, 5.8, 0.4, size=14, bold=True, color=GREEN)

txt(slide, "Ou passo a passo:", 6.7, 2.9, 6.1, 0.38, size=13, color=LIGHT_TXT)
cmds = [
    "python duck_init.py",
    "python run_staging.py",
    "python run_journey.py",
    "python run_service_metrics.py",
    "python run_human_escalation_metrics.py",
    "python run_stakeholder_report.py",
]
y = 3.32
for cmd in cmds:
    add_rect(slide, 6.7, y, 6.1, 0.42, DARK_BG)
    txt(slide, cmd, 6.85, y+0.05, 5.8, 0.32, size=12, bold=True, color=GREEN)
    y += 0.47


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PRINCIPAIS ACHADOS / KPIs
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Principais Achados", "O que os dados revelaram após processamento")

kpis = [
    ("5.000", "Conversas\n(pós dedup)"),
    ("4.743", "Conversas\nResolvidas"),
    ("587", "Conversas\nReabertas"),
    ("1.587,94 min", "Tempo médio\nde Resolução"),
    ("14,77 min", "Tempo médio\nde Fila"),
    ("22,20 min", "Atendimento\nMédio IA"),
    ("434,84 min", "Atendimento\nMédio Humano"),
    ("1.813,58 min", "Maior resolução\n(Pacific Goods)"),
]

x = 0.35
y = 1.75
for i, (val, label) in enumerate(kpis):
    col = i % 4
    row = i // 4
    xl = 0.35 + col * 3.2
    yl = 1.75 + row * 2.0
    kpi_box(slide, label, val, xl, yl, w=3.05, h=1.75)

txt(slide, "N2 exige muito mais atendimento humano que N1:",
    0.4, 5.9, 6.0, 0.4, size=14, bold=True, color=TEAL)
txt(slide, "678,01 min (N2)  vs  255,99 min (N1)  em média",
    0.4, 6.35, 6.0, 0.4, size=14, color=LIGHT_TXT)

txt(slide, "Pacific Goods / Suporte tem a maior espera\nmédio após escalonamento humano:",
    6.8, 5.9, 6.1, 0.6, size=14, bold=True, color=TEAL)
txt(slide, "23,49 minutos", 6.8, 6.55, 6.1, 0.4, size=18, bold=True, color=YELLOW)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DECISOES TECNICAS
# ════════════════════════════════════════════════════════════════════════════
slide = content_slide("Decisões Técnicas")

decisoes = [
    ("DuckDB como engine",
     "Simples, rápido para CSVs locais, excelente SQL analítico sem infraestrutura externa."),
    ("SQL em arquivos separados por camada",
     "Lógica auditável: cada .sql corresponde a uma etapa, fácil de revisar e versionar."),
    ("SELECT DISTINCT na staging",
     "Remove duplicatas exatas sem precisar identificar a 'cópia correta' manualmente."),
    ("Chaves compostas como identificador único",
     "instance + account_id + conversation_id e instance + account_id + event_id."),
    ("Fallback de created_at",
     "394 conversas sem conversation_created → usa stg_conversations.created_at_utc."),
    ("Última resolução como resolution_time",
     "Representa o ciclo completo em conversas reabertas. Primeira resolução mantida separada."),
    ("Classificação de intervalos por estado",
     "queue, ai_service, human_service, resolved, terminal → cálculo direto por tipo."),
    ("Espera pós-escalonamento",
     "first_human_reply_created_at − first_human_assigned_at (campo consolidado, não evento granular)."),
]

y = 1.25
for i, (titulo, desc) in enumerate(decisoes):
    col = i % 2
    row = i // 2
    xl = 0.35 + col * 6.5
    yl = 1.25 + row * 1.5
    bg = ACCENT if col == 0 else MID_BG
    add_rect(slide, xl, yl, 6.1, 1.35, bg)
    txt(slide, titulo, xl+0.12, yl+0.08, 5.85, 0.4, size=13, bold=True, color=TEAL)
    txt(slide, desc, xl+0.12, yl+0.5, 5.85, 0.75, size=12, color=LIGHT_TXT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ARQUITETURA EM PRODUCAO
# ════════════════════════════════════════════════════════════════════════════
slide = dark_slide("Arquitetura em Produção", "Como escalar essa solução para o ambiente real")

layers = [
    ("MongoDB\n/ Postgres", "Dados\noperacionais", ORANGE),
    ("S3\nBronze", "CSVs / Parquet\nparticionado", RGBColor(0xC7, 0x6B, 0x15)),
    ("Redshift\nRaw", "Cópia fiel\ndos dados brutos", RGBColor(0x02, 0x96, 0xB3)),
    ("dbt\n+ Staging", "Limpeza e\ntipagem versionada", TEAL),
    ("Redshift\nGold", "Marts prontos\npara consumo", GREEN),
    ("BI\n(Metabase...)", "Consulta tabelas\npré-calculadas", RGBColor(0x9C, 0x27, 0xB0)),
]

x = 0.3
for label, sub, color in layers:
    add_rect(slide, x, 1.7, 1.95, 1.5, color)
    txt(slide, label, x, 1.75, 1.95, 0.75, size=13, bold=True,
        color=DARK_BG if color in (TEAL, GREEN) else WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sub, x, 2.45, 1.95, 0.65, size=11,
        color=DARK_BG if color in (TEAL, GREEN) else WHITE, align=PP_ALIGN.CENTER)
    if x < 0.3 + 5 * 2.18:
        txt(slide, "→", x+1.95, 2.2, 0.25, 0.5, size=18, bold=True, color=LIGHT_TXT)
    x += 2.18

pilares = [
    ("Performance", [
        "Tabelas gold materializadas consultadas pelo BI",
        "Sort/distribution keys em Redshift (account_id, created_at)",
        "Jobs incrementais por data de evento",
    ], ORANGE),
    ("Visibilidade", [
        "Catálogo de métricas com definição, fórmula e owner",
        "KPIs por cliente, time, inbox, agente, N1/N2, IA vs humano",
        "Checks de qualidade publicados e monitorados",
    ], TEAL),
    ("Regras de Negócio", [
        "SQL/dbt versionado: uma definição única por métrica",
        "Conversas reabertas: primeira e última resolução",
        "Alertas para divergência, queda de volume, nulos",
    ], GREEN),
    ("Orquestração", [
        "Airflow / Dagster / Prefect para execução",
        "dbt para transformação, testes e documentação",
        "Logs e alertas no CloudWatch ou equivalente",
    ], RGBColor(0x9C, 0x27, 0xB0)),
]

x = 0.3
for pilar, items, color in pilares:
    add_rect(slide, x, 3.45, 2.9, 2.95, ACCENT)
    add_rect(slide, x, 3.45, 2.9, 0.45, color)
    txt(slide, pilar, x+0.1, 3.48, 2.7, 0.4, size=14, bold=True,
        color=DARK_BG if color in (TEAL, GREEN) else WHITE, align=PP_ALIGN.CENTER)
    y_item = 4.0
    for item in items:
        txt(slide, "• " + item, x+0.12, y_item, 2.65, 0.6, size=11, color=LIGHT_TXT)
        y_item += 0.7
    x += 3.1


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CONCLUSAO
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
add_rect(slide, 0, 0, 0.15, 7.5, TEAL)
add_rect(slide, 0.15, 0, 13.18, 0.08, TEAL)
add_rect(slide, 0.15, 7.42, 13.18, 0.08, TEAL)

txt(slide, "Conclusão", 0.5, 0.4, 12.5, 0.65, size=30, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 0.5, 1.2, 12.3, 1.05, ACCENT)
txt(slide,
    "O problema central não é lentidão — é uma arquitetura onde o BI calcula a métrica\n"
    "em vez de consumir a métrica.",
    0.65, 1.28, 12.0, 0.88, size=16, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

pontos = [
    ("O que foi entregue neste case",
     "Pipeline local completo em DuckDB simulando a arquitetura em camadas: ingestão, staging, "
     "jornada, métricas de serviço, escalonamento humano e relatório visual.",
     TEAL),
    ("Impacto em Produção",
     "Dashboards mais rápidos (consultam tabelas pré-calculadas), métricas centralizadas e "
     "versionadas, visibilidade total sobre fila, IA, humano e escalonamento.",
     GREEN),
    ("Próximos Passos Recomendados",
     "Migrar modelagem para dbt + Redshift, orquestrar com Airflow/Dagster, publicar catálogo "
     "de métricas e ativar alertas de qualidade de dados.",
     ORANGE),
]

y = 2.45
for titulo, desc, color in pontos:
    add_rect(slide, 0.5, y, 12.3, 1.35, MID_BG)
    add_rect(slide, 0.5, y, 0.08, 1.35, color)
    txt(slide, titulo, 0.72, y+0.1, 12.0, 0.42, size=15, bold=True, color=color)
    txt(slide, desc, 0.72, y+0.55, 11.9, 0.7, size=13, color=LIGHT_TXT)
    y += 1.5

add_rect(slide, 0.5, 6.9, 12.3, 0.48, ACCENT)
txt(slide,
    "\"O dashboard não deveria calcular a métrica — ele deveria consumir a métrica.\"",
    0.65, 6.95, 12.0, 0.38, size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER,
    italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SALVAR
# ════════════════════════════════════════════════════════════════════════════
out = r"d:\Backup PC\Code\Estudos\case_ch\Cloud_Humans_DataEngineering_Case.pptx"
prs.save(out)
print(f"Salvo em: {out}")